from spider import json_load
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import pandas as pd
from numpy import nan


class PPT:
    def __init__(self, input_file='template.pptx'):
        self.prs = Presentation(input_file)
        self.slides = self.prs.slides

    def save(self, output_file='output.pptx'):
        self.prs.save(output_file)

    def analyze_layouts(self, outputfile='output.pptx'):
        # 遍历每个版式与占位符
        for s, layout in enumerate(self.prs.slide_layouts):
            slide = self.prs.slides.add_slide(layout)

            # 是否有标题占位符
            try:
                title = slide.shapes.title
                title.text = f'{s}样式-标题'
            except AttributeError:
                print(f'>>> page {s} has no title')

            # 将其他占位符(placeholders)命名为x样式x号
            for each in slide.placeholders:
                each.text = f'{s}样式-{each.placeholder_format.idx}号'

            # 保存
            self.save(outputfile)

    def chart(self, df, placehold):
        data = ChartData()
        # 横轴名
        data.categories = list(df.index)
        # 系列名，数据
        for col_name in df:
            array = df[col_name].tolist()
            data.add_series(col_name, array)
        # 图表类型
        type = XL_CHART_TYPE.COLUMN_CLUSTERED
        # 指定placehold中插入图表
        placehold.insert_chart(type, data)


class CityPPT(PPT):
    path = r'E:\city_report\合肥\2017-09-06'

    def excel2chart(self, placehold, file, sheet, index, columns=[]):
        df = pd.read_excel(file, sheet, index_col=0)
        df = df.replace(nan, 0)
        df.index = index
        if columns:
            df.columns = columns
        self.chart(df, placehold)

    def new_slide(self, layout_index):
        layouts = self.prs.slide_layouts
        slide = self.prs.slides.add_slide(layouts[layout_index])
        placehold = slide.placeholders
        return slide, placehold

    def new_plate_slide(self, plate):
        # add a new slide
        slide, placehold = self.new_slide(1)

        # 板块名
        placehold[11].text = plate

        file = f'{self.path}/{plate}土地.xlsx'
        index = ['2013', '2014', '2015', '2016', '2017.1-8']

        # chart1
        sheet = '土地成交结构'
        self.excel2chart(placehold[14], file, sheet, index)

        # chart2
        sheet = '土地住宅面积'
        columns = ['可建面(万㎡)', '住宅面积（万㎡）']
        self.excel2chart(placehold[15], file, sheet, index, columns)

        # chart3
        sheet = '土地价格'
        columns = ['楼面价(元/㎡)', '地价房价比']
        self.excel2chart(placehold[16], file, sheet, index, columns)

    def new_pianqu_slide(self, plate, pianqu):
        # add a new slide
        slide, placehold = self.new_slide(2)

        # 删除片区名称中多余的字
        pianqu_ = pianqu.replace('板块', '').replace('片区', '')

        # 板块名及标题
        placehold[11].text = plate
        placehold[12].text = f'【{pianqu_}-市场量价】'

        file = f'{self.path}/{plate}_{pianqu}.xlsx'
        index_year = list(f'{x:02d}' for x in range(7, 17))
        index_year.append('17.1-8')
        index_year_ = list(range(2007, 2017))
        index_month = list(f'17.0{x}' for x in range(1, 9))
        columns_gxj = ['上市面积(万㎡)', '成交面积(万㎡)', '均价(元/㎡)']
        columns_stk = ['库存(万㎡)', '去化周期(月)']

        # chart1
        placehold[14].text = '2007-2017.8{pianqu_}商品住宅年度供销走势'
        sheet = '年度供销价'
        self.excel2chart(placehold[15], file, sheet, index_year, columns_gxj)

        # chart2
        placehold[16].text = '2017年1-8月{pianqu_}商品住宅月度供销走势'
        sheet = '月度供销价'
        self.excel2chart(placehold[17], file, sheet, index_month, columns_gxj)

        # chart3
        placehold[18].text = '2007-2016{pianqu_}商品住宅存量及去化周期年度走势'
        sheet = '年度库存'
        self.excel2chart(placehold[19], file, sheet, index_year_, columns_stk)

        # chart4
        placehold[20].text = '2017.1-2017.8{pianqu_}商品住宅存量及去化周期月度走势'
        sheet = '月度库存'
        self.excel2chart(placehold[21], file, sheet, index_month, columns_stk)


if __name__ == '__main__':
    area_dict = json_load('area_dict')

    for plate in area_dict:
        print('=' * 20, plate, '=' * 20)
        # 创建一个新的ppt
        ppt = CityPPT()

        # 板块土地
        ppt.new_plate_slide(plate)

        # 各片区市场
        for pianqu in area_dict[plate]:
            print(f'>>> {pianqu}')
            ppt.new_pianqu_slide(plate, pianqu)

        # 保存
        print('>>> saving...')
        ppt.save(f'{ppt.path}/{plate}.pptx')
        print('>>> ok!')
