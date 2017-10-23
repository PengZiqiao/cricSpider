import xlwings as xw
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE


class ToExcel:
    """将df存入Excel"""

    def __init__(self):
        self.app = xw.App(visible=False, add_book=False)

    def creat_book(self, path):
        """创建一个新的excel文件"""
        self.wb = self.app.books.add()
        self.wb.save(path)


    def df2sheet(self, sheet, df):
        """在指定excel文件下创建一个新的sheet,并写上df保存"""
        sheets = self.wb.sheets
        sht = sheets.add(sheet, after=len(sheets))
        sht.range('A1').value = df
        self.wb.save()
        print('>>> 保存成功！')

    def rewrite_sheet(self, book, sheet, df):
        wb = self.app.books.open(book)
        sht = wb.sheets[sheet]
        sht.clear()
        sht.range('A1').value = df
        wb.save()
        wb.close()
        print('>>> 保存成功！')

    def close(self):
        self.wb.close()


class PPT:
    def __init__(self, input_file='template.pptx'):
        self.prs = Presentation(input_file)
        self.slides = self.prs.slides

    def save(self, output_file='output.pptx'):
        self.prs.save(output_file)

    def analyze_layouts(self, outputfile='output.pptx'):
        # 遍历每个版式与占位符
        for s, layout in enumerate(self.prs.slide_layouts):
            slide = self.prs.slides.add_slide(layout)

            # 是否有标题占位符
            try:
                title = slide.shapes.title
                title.text = f'{s}样式-标题'
            except AttributeError:
                print(f'>>> page {s} has no title')

            # 将其他占位符(placeholders)命名为x样式x号
            for each in slide.placeholders:
                each.text = f'{s}样式-{each.placeholder_format.idx}号'

            # 保存
            self.save(outputfile)

    def chart(self, df, placehold):
        data = ChartData()
        # 横轴名
        data.categories = list(df.index)
        # 系列名，数据
        for col_name in df:
            array = df[col_name].tolist()
            data.add_series(col_name, array)
        # 图表类型
        type = XL_CHART_TYPE.COLUMN_CLUSTERED
        # 指定placehold中插入图表
        placehold.insert_chart(type, data)
